"""
Multi-format content extractor.
Supports: PDF, DOCX, PPTX, Image (OCR), Markdown, Text, Web URL.
"""
import os
import re
from pathlib import Path
from typing import Optional

from backend.models.schemas import ExtractionResult, SourceType


def detect_source_type(filename: str) -> SourceType:
    """Detect source type from file extension."""
    ext = Path(filename).suffix.lower()
    mapping = {
        ".pdf": SourceType.PDF,
        ".docx": SourceType.DOCX,
        ".doc": SourceType.DOCX,
        ".pptx": SourceType.PPTX,
        ".ppt": SourceType.PPTX,
        ".png": SourceType.IMAGE,
        ".jpg": SourceType.IMAGE,
        ".jpeg": SourceType.IMAGE,
        ".gif": SourceType.IMAGE,
        ".bmp": SourceType.IMAGE,
        ".webp": SourceType.IMAGE,
        ".md": SourceType.MARKDOWN,
        ".txt": SourceType.TEXT,
        ".csv": SourceType.TEXT,
    }
    return mapping.get(ext, SourceType.TEXT)


async def extract_content(file_path: str, source_type: SourceType) -> ExtractionResult:
    """Extract text content from a file based on its type."""
    extractors = {
        SourceType.PDF: _extract_pdf,
        SourceType.DOCX: _extract_docx,
        SourceType.PPTX: _extract_pptx,
        SourceType.IMAGE: _extract_image,
        SourceType.MARKDOWN: _extract_markdown,
        SourceType.TEXT: _extract_text,
    }

    extractor = extractors.get(source_type, _extract_text)
    return await extractor(file_path)


async def extract_from_url(url: str) -> ExtractionResult:
    """Extract content from a web URL using trafilatura with a custom User-Agent Header fallback."""
    try:
        from trafilatura import fetch_url, extract, bare_extraction
        import urllib.request
        from urllib.error import URLError, HTTPError

        downloaded = None
        # Use Standard Chrome User-Agent to evade basic scraper walls
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
            'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
            'Accept-Language': 'ko-KR,ko;q=0.9,en-US;q=0.8,en;q=0.7',
        }
        
        try:
            req = urllib.request.Request(url, headers=headers)
            with urllib.request.urlopen(req, timeout=10) as response:
                downloaded = response.read().decode('utf-8', errors='ignore')
        except Exception as e:
            print(f"urllib Request failed for {url}: {e}")
            # Fallback to pure trafilatura fetch
            downloaded = fetch_url(url)

        if not downloaded:
            return ExtractionResult(
                text="",
                title=url,
                metadata={"error": "Failed to fetch URL. Blocked by server or Timeout.", "url": url},
            )

        # Get full extraction with metadata
        result = bare_extraction(downloaded, include_comments=False)

        text = result.get("text", "") if result else ""
        title = result.get("title", url) if result else url
        
        # Simple extract fallback if bare_extraction returned empty text
        if not text.strip():
            text = extract(downloaded, include_comments=False) or ""

        metadata = {
            "url": url,
            "author": result.get("author", "") if result else "",
            "date": result.get("date", "") if result else "",
            "sitename": result.get("sitename", "") if result else "",
            "description": result.get("description", "") if result else "",
        }

        return ExtractionResult(text=text, title=title, metadata=metadata)

    except ImportError:
        return ExtractionResult(
            text="",
            title=url,
            metadata={"error": "trafilatura not installed", "url": url},
        )
    except Exception as e:
        return ExtractionResult(
            text="",
            title=url,
            metadata={"error": str(e), "url": url},
        )


async def _extract_pdf(file_path: str) -> ExtractionResult:
    """Extract text from PDF using PyMuPDF."""
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(file_path)
        text_parts = []
        for page in doc:
            text_parts.append(page.get_text())
        doc.close()

        full_text = "\n\n".join(text_parts)
        title = _extract_title_from_text(full_text) or Path(file_path).stem

        return ExtractionResult(
            text=full_text,
            title=title,
            metadata={"pages": len(text_parts), "source_file": os.path.basename(file_path)},
        )
    except Exception as e:
        return ExtractionResult(
            text="",
            title=Path(file_path).stem,
            metadata={"error": str(e)},
        )


async def _extract_docx(file_path: str) -> ExtractionResult:
    """Extract text from DOCX using python-docx."""
    try:
        from docx import Document

        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        full_text = "\n\n".join(paragraphs)
        title = paragraphs[0] if paragraphs else Path(file_path).stem

        return ExtractionResult(
            text=full_text,
            title=title,
            metadata={
                "paragraphs": len(paragraphs),
                "source_file": os.path.basename(file_path),
            },
        )
    except Exception as e:
        return ExtractionResult(
            text="",
            title=Path(file_path).stem,
            metadata={"error": str(e)},
        )


async def _extract_pptx(file_path: str) -> ExtractionResult:
    """Extract text from PPTX using python-pptx."""
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        slides_text = []

        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                slides_text.append(f"[Slide {i}]\n" + "\n".join(slide_texts))

        full_text = "\n\n".join(slides_text)
        title = _extract_title_from_text(full_text) or Path(file_path).stem

        return ExtractionResult(
            text=full_text,
            title=title,
            metadata={
                "slides": len(prs.slides),
                "source_file": os.path.basename(file_path),
            },
        )
    except Exception as e:
        return ExtractionResult(
            text="",
            title=Path(file_path).stem,
            metadata={"error": str(e)},
        )


async def _extract_image(file_path: str) -> ExtractionResult:
    """Extract text from image using OCR (pytesseract)."""
    try:
        import pytesseract
        from PIL import Image

        img = Image.open(file_path)
        # Use Korean + English for OCR
        text = pytesseract.image_to_string(img, lang="kor+eng")

        return ExtractionResult(
            text=text.strip(),
            title=Path(file_path).stem,
            metadata={
                "source_file": os.path.basename(file_path),
                "image_size": f"{img.width}x{img.height}",
                "ocr": True,
            },
        )
    except ImportError:
        return ExtractionResult(
            text="[이미지 파일 - OCR 미설치 (pytesseract)]",
            title=Path(file_path).stem,
            metadata={
                "source_file": os.path.basename(file_path),
                "error": "pytesseract not available",
            },
        )
    except Exception as e:
        return ExtractionResult(
            text="",
            title=Path(file_path).stem,
            metadata={"error": str(e)},
        )


async def _extract_markdown(file_path: str) -> ExtractionResult:
    """Extract content from Markdown file."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()

        # Extract title from first heading
        title_match = re.search(r"^#\s+(.+)$", text, re.MULTILINE)
        title = title_match.group(1) if title_match else Path(file_path).stem

        return ExtractionResult(
            text=text,
            title=title,
            metadata={"source_file": os.path.basename(file_path), "format": "markdown"},
        )
    except Exception as e:
        return ExtractionResult(
            text="",
            title=Path(file_path).stem,
            metadata={"error": str(e)},
        )


async def _extract_text(file_path: str) -> ExtractionResult:
    """Extract content from plain text file."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()

        title = _extract_title_from_text(text) or Path(file_path).stem

        return ExtractionResult(
            text=text,
            title=title,
            metadata={"source_file": os.path.basename(file_path)},
        )
    except UnicodeDecodeError:
        # Try with different encoding
        try:
            with open(file_path, "r", encoding="cp949") as f:
                text = f.read()
            return ExtractionResult(
                text=text,
                title=Path(file_path).stem,
                metadata={"source_file": os.path.basename(file_path), "encoding": "cp949"},
            )
        except Exception as e:
            return ExtractionResult(
                text="",
                title=Path(file_path).stem,
                metadata={"error": str(e)},
            )
    except Exception as e:
        return ExtractionResult(
            text="",
            title=Path(file_path).stem,
            metadata={"error": str(e)},
        )


def _extract_title_from_text(text: str) -> Optional[str]:
    """Try to extract a title from the first line of text."""
    if not text:
        return None
    lines = text.strip().split("\n")
    for line in lines[:5]:
        cleaned = line.strip().lstrip("#").strip()
        if cleaned and len(cleaned) > 2 and len(cleaned) < 200:
            return cleaned
    return None
